from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from school_assistant.file_text import extract_local_file


class FileTextTests(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)

    def tearDown(self):
        self.temp.cleanup()

    def test_plain_docx_xlsx_and_pptx(self):
        text_path = self.root / "通知.txt"
        text_path.write_text("请在周五前提交登记表", encoding="utf-8")

        docx_path = self.root / "通知.docx"
        document = Document()
        document.add_paragraph("9月1日上午报到")
        document.save(docx_path)

        xlsx_path = self.root / "名单.xlsx"
        workbook = Workbook()
        workbook.active.append(["姓名", "截止时间"])
        workbook.active.append(["张三", "2026-09-02"])
        workbook.save(xlsx_path)
        workbook.close()

        pptx_path = self.root / "安排.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "开学会议安排"
        presentation.save(pptx_path)

        expectations = {
            text_path: "提交登记表",
            docx_path: "上午报到",
            xlsx_path: "截止时间",
            pptx_path: "开学会议安排",
        }
        for path, expected in expectations.items():
            with self.subTest(path=path.suffix):
                result = extract_local_file(str(path))
                self.assertEqual(result["state"], "extracted")
                self.assertIn(expected, result["text"])

    def test_scanned_like_pdf_and_unsupported_file_need_manual_review(self):
        pdf_path = self.root / "扫描件.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=100, height=100)
        with pdf_path.open("wb") as stream:
            writer.write(stream)
        self.assertEqual(extract_local_file(str(pdf_path))["state"], "empty")

        binary_path = self.root / "旧格式.doc"
        binary_path.write_bytes(b"dummy")
        self.assertEqual(extract_local_file(str(binary_path))["state"], "unsupported")


if __name__ == "__main__":
    unittest.main()
