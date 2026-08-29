"""Extract text from downloaded attachments without OCR or Office automation."""

from __future__ import annotations

import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Iterable


MAX_FILE_BYTES = 50 * 1024 * 1024
MAX_TEXT_CHARS = 30_000
PLAIN_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".json", ".jsonl", ".xml", ".log", ".ini", ".yaml", ".yml",
}
SUPPORTED_EXTENSIONS = PLAIN_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".xlsm", ".xls", ".pptx", ".html", ".htm"}


class _HTMLText(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(data.strip())


def _bounded(parts: Iterable[Any], limit: int) -> str:
    result: list[str] = []
    size = 0
    for raw in parts:
        value = re.sub(r"[ \t]+", " ", str(raw or "")).strip()
        if not value:
            continue
        remaining = limit - size
        if remaining <= 0:
            break
        result.append(value[:remaining])
        size += min(len(value), remaining) + 1
    return "\n".join(result).strip()


def _read_plain(path: Path, limit: int) -> str:
    raw = path.read_bytes()[: min(path.stat().st_size, limit * 4 + 4096)]
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")
    return text[:limit].strip()


def _read_pdf(path: Path, limit: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path), strict=False)
    return _bounded((page.extract_text() or "" for page in reader.pages), limit)


def _read_docx(path: Path, limit: int) -> str:
    from docx import Document

    document = Document(str(path))
    parts: list[str] = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return _bounded(parts, limit)


def _read_xlsx(path: Path, limit: int) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True, keep_links=False)
    try:
        parts: list[str] = []
        size = 0
        for sheet in workbook.worksheets:
            parts.append(f"[工作表：{sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                line = " | ".join(str(value) for value in row if value not in (None, ""))
                if line:
                    parts.append(line)
                    size += len(line) + 1
                if size >= limit:
                    break
            if size >= limit:
                break
        return _bounded(parts, limit)
    finally:
        workbook.close()


def _read_xls(path: Path, limit: int) -> str:
    import xlrd

    workbook = xlrd.open_workbook(str(path), on_demand=True)
    try:
        parts: list[str] = []
        size = 0
        for sheet in workbook.sheets():
            parts.append(f"[工作表：{sheet.name}]")
            for row_number in range(sheet.nrows):
                line = " | ".join(str(value) for value in sheet.row_values(row_number) if value not in (None, ""))
                if line:
                    parts.append(line)
                    size += len(line) + 1
                if size >= limit:
                    break
            if size >= limit:
                break
        return _bounded(parts, limit)
    finally:
        workbook.release_resources()


def _read_pptx(path: Path, limit: int) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        parts.append(f"[第 {index} 页]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
    return _bounded(parts, limit)


def _read_html(path: Path, limit: int) -> str:
    parser = _HTMLText()
    parser.feed(_read_plain(path, limit * 2))
    return _bounded(parser.parts, limit)


def extract_local_file(path_value: str, limit: int = MAX_TEXT_CHARS) -> dict[str, str]:
    """Return state/text/error. Never runs macros, Office, JavaScript, or OCR."""
    path = Path(path_value)
    if not path.is_file():
        return {"state": "error", "text": "", "error": "文件不存在"}
    try:
        size = path.stat().st_size
        if size > MAX_FILE_BYTES:
            return {"state": "unsupported", "text": "", "error": "文件超过 50MB，未自动读取"}
        extension = path.suffix.lower()
        if extension not in SUPPORTED_EXTENSIONS:
            return {"state": "unsupported", "text": "", "error": f"暂不支持 {extension or '无扩展名'} 格式"}
        if extension in PLAIN_EXTENSIONS:
            text = _read_plain(path, limit)
        elif extension == ".pdf":
            text = _read_pdf(path, limit)
        elif extension == ".docx":
            text = _read_docx(path, limit)
        elif extension in {".xlsx", ".xlsm"}:
            text = _read_xlsx(path, limit)
        elif extension == ".xls":
            text = _read_xls(path, limit)
        elif extension == ".pptx":
            text = _read_pptx(path, limit)
        else:
            text = _read_html(path, limit)
        if not text.strip():
            return {"state": "empty", "text": "", "error": "未提取到文本（可能是扫描件；未启用 OCR）"}
        return {"state": "extracted", "text": text[:limit], "error": ""}
    except Exception as exc:
        return {"state": "error", "text": "", "error": f"{type(exc).__name__}: {exc}"[:300]}
